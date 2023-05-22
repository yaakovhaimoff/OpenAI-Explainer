import argparse
import asyncio
import time
import backoff
from pptx import Presentation
import openai


async def process_presentation() -> list[str]:
    """
    Process each slide in the presentation asynchronously and retrieve the responses.

    Returns:
        list[str]: List of explanations for each slide.
    """
    prs = Presentation(presentation_path)
    tasks = []
    explanations = []
    for slide in prs.slides:
        tasks.append(asyncio.create_task(process_slide(slide, explanations)))
    await asyncio.gather(*tasks)
    return explanations


@backoff.on_exception(backoff.expo, openai.error.RateLimitError)
async def process_slide(slide, explanations) -> None:
    """
    Process a slide and append the response to the explanations list.
    If an error occurs, save an informative error message.

    Args:
        slide: Slide object from the PowerPoint presentation.
        explanations (list[str]): List to store the explanations.
    """
    try:
        response = await get_response(slide)
        explanations.append(response)
    except Exception as e:
        error_message = f"Error processing slide: {str(e)}"
        explanations.append(error_message)


@backoff.on_exception(backoff.expo, openai.error.RateLimitError)
async def get_response(slide) -> str:
    """
    Send a request to the OpenAI Chat API to get a response for the given slide.

    Args:
        slide: Slide object from the PowerPoint presentation.

    Returns:
        str: Response generated by the OpenAI Chat API.
    """
    messages.append({"role": "user", "content": process_slide_text(slide)})
    response = await openai.ChatCompletion.acreate(
        model="gpt-3.5-turbo",
        messages=messages
    )
    return response["choices"][0].message.content


def process_slide_text(slide) -> str:
    """
    Extract the text content from a slide in the PowerPoint presentation.

    Args:
        slide: Slide object from the PowerPoint presentation.

    Returns:
        str: Processed text content from the slide.
    """
    slide_text = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                slide_text.append(run.text.strip())
    return " ".join(slide_text)


def print_to_file(explanations) -> None:
    """
    Print the explanations to a JSON file.

    Args:
        explanations (list[str]): List of explanations to be printed.
    """
    output_file = f'{presentation_path.split("/")[-1].split(".")[0]}.json'
    with open(output_file, 'w') as f:
        for i, explanation in enumerate(explanations):
            f.write(f'chat GPT answer {i + 1} {str(explanation)} \n\n')


async def main() -> None:
    """
    Main function to run the program.
    """
    openai.api_key = "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"
    explanations = await process_presentation()
    print_to_file(explanations)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Process PowerPoint presentation and generate explanations.')
    parser.add_argument('presentation', metavar='presentation', type=str,
                        help='path to the PowerPoint presentation')
    args = parser.parse_args()

    messages = [
        {"role": "system", "content": "Please summarize the slides and provide additional information."}
    ]

    presentation_path = args.presentation
    # presentation_path = "final.pptx"
    print(f"Processing {presentation_path}")
    start_time = time.time()
    asyncio.run(main())

    end_time = time.time()
    execution_time = end_time - start_time
    minutes, seconds = divmod(execution_time, 60)
    print(f"Execution time: {minutes:.0f} minutes {seconds:.2f} seconds")